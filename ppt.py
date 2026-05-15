import os
import re
from typing import Optional, List, Dict
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image
import io
import tempfile

# Novos módulos para SharePoint e XML
from sharepoint_client import SharePointClient
from data_manager import DataManager
from pdf_parser import extrair_produtos_do_pdf, validar_produtos

# Configurações do novo sistema
xml_orcamento = r"C:\Users\pedro\OneDrive\Desktop\lucenera\orcamento.xml"  # Será definido dinamicamente
excel_master = r"C:\Users\pedro\OneDrive\Desktop\lucenera\master_produtos.xlsx"  # Caminho do Excel master
logo_path = r"C:\script python\script python power point\logo.png"

# Instanciar clientes
sharepoint_client = SharePointClient()
data_manager = None  # Será inicializado com dados reais

# Extrair nome do XML sem extensão
xml_nome = "orcamento_sharepoint"  # Nome padrão, será atualizado dinamicamente
ppt_saida = rf"C:\Users\pedro\OneDrive\Desktop\lucenera\Ficha_Tecnica_{xml_nome}.pptx"

img_exts = {'.jpg', '.jpeg', '.png', '.gif', '.bmp'}

a4_width = Inches(8.27)
a4_height = Inches(11.69)

# Função para processar orçamento XML e gerar PowerPoint
def gerar_powerpoint_sharepoint(xml_path: str, excel_path: str, output_path: str = None):
    """
    Gera PowerPoint a partir de XML usando dados do SharePoint
    Agora com agrupamento por código L (um slide por grupo)
    
    Args:
        xml_path: Caminho para arquivo XML de orçamento
        excel_path: Caminho para Excel master de produtos
        output_path: Caminho de saída (opcional)
    """
    try:
        print("[INICIO] Iniciando geração PowerPoint com SharePoint...")
        
        # Inicializar gerenciador de dados
        global data_manager
        data_manager = DataManager(xml_path, excel_path)
        
        # Processar orçamento (agora retorna GRUPOS por código L)
        print("[XML] Processando orçamento XML...")
        grupos = data_manager.processar_orcamento()
        
        if not grupos:
            raise Exception("Nenhum grupo encontrado no XML")
        
        print(f"[OK] {len(grupos)} grupos encontrados (agrupados por código L)!")
        
        # Definir nome de saída
        if output_path:
            ppt_saida_final = output_path
        else:
            xml_base = os.path.splitext(os.path.basename(xml_path))[0]
            ppt_saida_final = rf"C:\Users\pedro\OneDrive\Desktop\lucenera\Ficha_Tecnica_{xml_base}.pptx"
        
        # Criar apresentação
        prs = Presentation()
        prs.slide_width = a4_width
        prs.slide_height = a4_height
        
        # Slide de capa
        criar_slide_capa(prs)
        
        # Processar cada GRUPO (código L)
        for i, grupo in enumerate(grupos, 1):
            lnum = grupo.get('lnum', '?')
            mandante = grupo.get('mandante')
            componentes = grupo.get('componentes', [])
            total = grupo.get('total_produtos', 1)
            
            print(f"\n[GRUPO {i}/{len(grupos)}] L{lnum} - {total} produto(s)")
            
            if mandante:
                print(f"  [MANDANTE] {mandante.get('codigo', '?')} - {mandante.get('categoria', 'SEM CATEGORIA')}")
            
            if componentes:
                print(f"  [COMPONENTES] {len(componentes)} item(ns):")
                for comp in componentes:
                    print(f"    - {comp.get('codigo', '?')}: {comp.get('categoria', 'SEM CATEGORIA')}")
            
            # Criar slide para o GRUPO (usando mandante como base)
            criar_slide_grupo(prs, grupo, sharepoint_client)
        
        # Salvar PowerPoint
        prs.save(ppt_saida_final)
        
        print(f"\n[SUCESSO] PowerPoint gerado com sucesso!")
        print(f"[ARQUIVO] {ppt_saida_final}")
        print(f"[TOTAL] {len(grupos)} grupos processados (1 slide por grupo)")
        return ppt_saida_final
        
    except Exception as e:
        print(f"[ERRO] Erro ao gerar PowerPoint: {str(e)}")
        raise


def gerar_powerpoint_pdf(pdf_path: str, excel_path: str, output_path: str = None):
    """
    Gera PowerPoint a partir de PDF de orçamento usando SharePoint
    
    Args:
        pdf_path: Caminho para arquivo PDF de orçamento
        excel_path: Caminho para Excel master de produtos  
        output_path: Caminho de saída (opcional)
    
    Returns:
        Caminho do arquivo PowerPoint gerado
    """
    try:
        print("[INICIO] Iniciando geração PowerPoint a partir de PDF...")
        
        # Extrair produtos do PDF
        print("[PDF] Extraindo produtos do PDF...")
        produtos_pdf = extrair_produtos_do_pdf(pdf_path)
        validar_produtos(produtos_pdf)
        
        print(f"[OK] {len(produtos_pdf)} produtos extraídos do PDF!")
        
        # Carregar Excel master para enriquecer dados
        print("[EXCEL] Carregando dados do Excel master...")
        from data_manager import ExcelMaster
        excel_master_obj = ExcelMaster(excel_path)
        
        # Enriquecer produtos com dados do Excel
        produtos_enriquecidos = []
        for p in produtos_pdf:
            codigo = p['codigo']
            
            # Buscar marca no Excel (se não veio do PDF)
            if p.get('marca'):
                marca = p['marca']  # Usar marca extraída do PDF
            else:
                marca = excel_master_obj.buscar_marca_por_codigo(codigo)
                if not marca:
                    marca = 'LUCENERA'  # Default
            
            # Criar produto enriquecido
            produto = {
                'lnum': p['item'].replace('L', '').replace('Item ', '').zfill(2),
                'codigo': codigo,
                'ref': p.get('referencia', codigo),
                'marca': marca,
                'quantidade': p.get('quantidade', '1'),
                'unidade': p.get('unidade', 'UN'),
                'descricao': p.get('descricao', '')
            }
            
            produtos_enriquecidos.append(produto)
            print(f"  [PRODUTO] {produto['lnum']}: {marca} - {codigo} - Ref: {produto['ref']} ({produto['quantidade']} {produto['unidade']})")
        
        # Definir nome de saída
        if output_path:
            ppt_saida_final = output_path
        else:
            pdf_base = os.path.splitext(os.path.basename(pdf_path))[0]
            ppt_saida_final = rf"C:\Users\pedro\OneDrive\Desktop\lucenera\Ficha_Tecnica_{pdf_base}.pptx"
        
        # Criar apresentação
        prs = Presentation()
        prs.slide_width = a4_width
        prs.slide_height = a4_height
        
        # Slide de capa
        criar_slide_capa(prs)
        
        # Processar cada produto
        sp_client = SharePointClient()
        for i, produto in enumerate(produtos_enriquecidos, 1):
            print(f"[PROC] Processando produto {i}/{len(produtos_enriquecidos)}: L{produto.get('lnum', '?')} - {produto.get('codigo', '?')}")
            
            # Criar slides para o produto
            criar_slides_produto(prs, produto, sp_client)
        
        # Salvar PowerPoint
        prs.save(ppt_saida_final)
        
        print(f"[SUCESSO] PowerPoint gerado com sucesso a partir de PDF!")
        print(f"[ARQUIVO] {ppt_saida_final}")
        print(f"[TOTAL] {len(produtos_enriquecidos)} produtos processados")
        return ppt_saida_final
        
    except Exception as e:
        print(f"[ERRO] Erro ao gerar PowerPoint de PDF: {str(e)}")
        raise


def criar_slide_capa(prs):
    """Cria slide de capa com logo Lucenera"""
    slide_capa = prs.slides.add_slide(prs.slide_layouts[6])
    if os.path.exists(logo_path):
        # Logo centralizada no meio da página
        logo_width = 3
        logo_height = 3
        logo_left = (8.27 - logo_width) / 2
        logo_top = (11.69 - logo_height) / 2
        slide_capa.shapes.add_picture(logo_path, Inches(logo_left), Inches(logo_top), width=Inches(logo_width))


def criar_slide_grupo(prs, grupo, sp_client):
    """
    Cria UM slide para um grupo completo (código L)
    Usa o item MANDANTE para busca no SharePoint
    Concatena componentes como texto adicional
    
    Args:
        prs: Apresentação PowerPoint
        grupo: Dicionário com {'lnum', 'mandante', 'componentes', 'total_produtos'}
        sp_client: Cliente SharePoint
    """
    lnum = grupo.get('lnum', '1')
    mandante = grupo.get('mandante')
    componentes = grupo.get('componentes', [])
    
    if not mandante:
        print(f"  [ERRO] Grupo L{lnum} sem item mandante!")
        return
    
    # Dados do mandante
    codigo_mandante = mandante.get('codigo', '')
    ref_mandante = mandante.get('ref', '')
    marca_mandante = mandante.get('marca', 'Interlight')
    categoria_mandante = mandante.get('categoria', '')
    
    print(f"  [SLIDE] Criando slide para grupo L{lnum}")
    print(f"    [MANDANTE] {codigo_mandante} - {categoria_mandante}")
    
    # Buscar arquivos no SharePoint pelo código do MANDANTE
    print(f"  [BUSCA_SHAREPOINT] Código: {codigo_mandante}")
    arquivos = sp_client.search_files_by_code(codigo_mandante)
    
    if not arquivos:
        print(f"  [AVISO] Nenhum arquivo encontrado para código {codigo_mandante}")
    else:
        print(f"  [OK] {len(arquivos)} arquivo(s) encontrado(s)")
    
    # Separar arquivos por tipo
    bula_files = [arq for arq in arquivos if arq.get('is_bula', False)]
    word_files = [arq for arq in arquivos if not arq.get('is_bula', False) and arq.get('type') == 'word']
    image_files = [arq for arq in arquivos if not arq.get('is_bula', False) and arq.get('type') == 'image']
    
    print(f"  [ARQUIVOS] {len(word_files)} ficha(s), {len(image_files)} imagem(ns), {len(bula_files)} bula(s)")
    
    # Criar slide de ficha técnica
    slide_ficha = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide_ficha)
    
    # Título pequeno em preto abaixo do cabeçalho
    tx_title = slide_ficha.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(3), Inches(0.4))
    tf_title = tx_title.text_frame
    tf_title.text = "LUMINÁRIA DE PROJETO"
    tf_title.paragraphs[0].font.name = 'Calibri'
    tf_title.paragraphs[0].font.size = Pt(12)
    tf_title.paragraphs[0].font.bold = False
    tf_title.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    
    # Número LXX grande em preto no canto superior direito
    tx_lnum = slide_ficha.shapes.add_textbox(Inches(6.8), Inches(1.2), Inches(1), Inches(0.5))
    tf_lnum = tx_lnum.text_frame
    tf_lnum.text = f"L{lnum}"
    tf_lnum.paragraphs[0].font.name = 'Calibri'
    tf_lnum.paragraphs[0].font.size = Pt(36)
    tf_lnum.paragraphs[0].font.bold = True
    tf_lnum.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    tf_lnum.paragraphs[0].alignment = PP_ALIGN.RIGHT
    
    # Linha contínua ANTES das fotos
    line1 = slide_ficha.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.5), Inches(2.0), Inches(7.77), Inches(2.0))
    line1.line.color.rgb = RGBColor(0, 0, 0)
    
    # Adicionar imagens do SharePoint (do mandante)
    adicionar_imagens_sharepoint(slide_ficha, word_files, image_files, sp_client, mandante)
    
    # Linha contínua DEPOIS das fotos
    line2 = slide_ficha.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.5), Inches(5.5), Inches(7.77), Inches(5.5))
    line2.line.color.rgb = RGBColor(0, 0, 0)
    
    # Adicionar especificações do SharePoint (incluindo componentes)
    adicionar_especificacoes_grupo(slide_ficha, word_files, sp_client, mandante, componentes)
    
    # Criar slide para CADA bula disponível (do mandante)
    if bula_files:
        print(f"  [BULA] {len(bula_files)} bula(s) encontrada(s)")
        for idx, bula_file in enumerate(bula_files, 1):
            print(f"  [BULA] Criando slide {idx}/{len(bula_files)}: {bula_file.get('name')}")
            criar_slide_bula(prs, bula_file, lnum, sp_client)


def criar_slides_produto(prs, produto, sp_client):
    """
    [FUNÇÃO LEGADA - Mantida para compatibilidade com sistema PDF]
    Cria slides para um produto específico usando busca SIMPLIFICADA por código interno
    
    Args:
        prs: Apresentação PowerPoint
        produto: Dados do produto
        sp_client: Cliente SharePoint
    """
    codigo = produto.get('codigo', '')
    ref = produto.get('ref', '')
    marca = produto.get('marca', 'Interlight')
    lnum = produto.get('lnum', '1')
    
    print(f"  [DEBUG PPT] Produto recebido - lnum: '{lnum}', codigo: '{codigo}', ref: '{ref}'")
    print(f"  [PROC] Processando produto L{lnum}: {marca} - Código: {codigo}")
    
    print(f"  [BUSCA_CODIGO] Buscando arquivos no SharePoint por código interno: {codigo}")
    
    # Buscar arquivos no SharePoint pelo CÓDIGO INTERNO (docx/jpg/png/pdf)
    arquivos = sp_client.search_files_by_code(codigo)
    
    if not arquivos:
        print(f"  [AVISO] Nenhum arquivo .docx encontrado para código {codigo}")
    else:
        print(f"  [OK] {len(arquivos)} arquivo(s) encontrado(s)")
        # Mostrar detalhes dos matches para debug
        for i, arq in enumerate(arquivos[:3]):  # Primeiro 3 apenas
            print(f"    [{i+1}] {arq['name']} (tipo: {arq.get('type', 'N/A')}, is_bula: {arq.get('is_bula', False)}, score: {arq.get('score', 0)})")
    
    # Separar arquivos por tipo: fichas vs bulas
    bula_files = [arq for arq in arquivos if arq.get('is_bula', False)]
    word_files = [arq for arq in arquivos if not arq.get('is_bula', False) and arq.get('type') == 'word']
    image_files = [arq for arq in arquivos if not arq.get('is_bula', False) and arq.get('type') == 'image']
    
    print(f"  [SEPARACAO] {len(word_files)} ficha(s), {len(image_files)} imagem(ns), {len(bula_files)} bula(s)")
    
    # Criar slide de ficha técnica
    slide_ficha = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide_ficha)
    
    # Título pequeno em preto abaixo do cabeçalho
    tx_title = slide_ficha.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(3), Inches(0.4))
    tf_title = tx_title.text_frame
    tf_title.text = "LUMINÁRIA DE PROJETO"
    tf_title.paragraphs[0].font.name = 'Calibri'
    tf_title.paragraphs[0].font.size = Pt(12)
    tf_title.paragraphs[0].font.bold = False
    tf_title.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    
    # Número LXX grande em preto no canto superior direito
    tx_lnum = slide_ficha.shapes.add_textbox(Inches(6.8), Inches(1.2), Inches(1), Inches(0.5))
    tf_lnum = tx_lnum.text_frame
    tf_lnum.text = f"L{lnum}"
    tf_lnum.paragraphs[0].font.name = 'Calibri'
    tf_lnum.paragraphs[0].font.size = Pt(36)
    tf_lnum.paragraphs[0].font.bold = True
    tf_lnum.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    tf_lnum.paragraphs[0].alignment = PP_ALIGN.RIGHT
    
    # Linha contínua ANTES das fotos
    line1 = slide_ficha.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.5), Inches(2.0), Inches(7.77), Inches(2.0))
    line1.line.color.rgb = RGBColor(0, 0, 0)
    
    # Adicionar imagens do SharePoint
    adicionar_imagens_sharepoint(slide_ficha, word_files, image_files, sp_client, produto)
    
    # Linha contínua DEPOIS das fotos
    line2 = slide_ficha.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.5), Inches(5.5), Inches(7.77), Inches(5.5))
    line2.line.color.rgb = RGBColor(0, 0, 0)
    
    # Adicionar especificações do SharePoint
    adicionar_especificacoes_sharepoint(slide_ficha, word_files, sp_client, produto)
    
    # Criar slide para CADA bula disponível (suporte a múltiplas bulas)
    if bula_files:
        print(f"  [BULA] {len(bula_files)} bula(s) encontrada(s)")
        for idx, bula_file in enumerate(bula_files, 1):
            print(f"  [BULA] Criando slide {idx}/{len(bula_files)}: {bula_file.get('name')}")
            criar_slide_bula(prs, bula_file, lnum, sp_client)


def adicionar_imagens_sharepoint(slide, word_files, image_files, sp_client, produto=None):
    """Adiciona imagens do SharePoint ao slide"""
    try:
        # Primeiro tentar extrair imagens do Word
        imagens_word = []
        if word_files:
            download_url = word_files[0]['download_url']
            imagens_word = sp_client.get_word_images(download_url)
        
        # Adicionar imagens do Word (produto e ambiente)
        imagens_adicionadas = 0
        for i, img_data in enumerate(imagens_word[:2]):
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as temp_img:
                temp_img.write(img_data)
                temp_path = temp_img.name
            
            try:
                if i == 0:
                    # Primeira imagem (produto) - menor à esquerda
                    slide.shapes.add_picture(temp_path, Inches(1.0), Inches(2.3), height=Inches(2.5))
                    imagens_adicionadas += 1
                elif i == 1:
                    # Segunda imagem (ambiente) - maior à direita
                    slide.shapes.add_picture(temp_path, Inches(4.2), Inches(2.3), height=Inches(3.0))
                    imagens_adicionadas += 1
            finally:
                os.unlink(temp_path)
        
        # Se não tiver no Word, usar imagens separadas
        if not imagens_word and image_files:
            for i, img_file in enumerate(image_files[:2]):
                try:
                    img_content = sp_client.download_file_content(img_file['download_url'])
                    
                    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as temp_img:
                        temp_img.write(img_content)
                        temp_path = temp_img.name
                    
                    if i == 0:
                        slide.shapes.add_picture(temp_path, Inches(1.0), Inches(2.3), height=Inches(2.5))
                        imagens_adicionadas += 1
                    elif i == 1:
                        slide.shapes.add_picture(temp_path, Inches(4.2), Inches(2.3), height=Inches(3.0))
                        imagens_adicionadas += 1
                        
                finally:
                    if os.path.exists(temp_path):
                        os.unlink(temp_path)
        
        # Se não adicionou nenhuma imagem e temos dados do produto, criar placeholders
        if imagens_adicionadas == 0 and produto:
            codigo = produto.get('codigo', '')
            ref = produto.get('ref', '') or produto.get('referencia', '')
            
            # Placeholder para imagem do produto (esquerda) - CÓDIGO DE REFERÊNCIA
            tx_img1 = slide.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(2.5), Inches(2.5))
            tf_img1 = tx_img1.text_frame
            tf_img1.text = f"REFERÊNCIA\n{ref if ref else 'N/A'}"
                
            tf_img1.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in tf_img1.paragraphs:
                p.font.name = 'Calibri'
                p.font.size = Pt(16)
                p.font.bold = True
                p.alignment = PP_ALIGN.CENTER
                
            # Borda para o placeholder
            tx_img1.line.color.rgb = RGBColor(100, 100, 100)
            tx_img1.line.width = Pt(1)
            tx_img1.fill.solid()
            tx_img1.fill.fore_color.rgb = RGBColor(245, 245, 245)
            
            # Placeholder para imagem de ambiente (direita) - CÓDIGO INTERNO
            tx_img2 = slide.shapes.add_textbox(Inches(4.2), Inches(2.3), Inches(2.5), Inches(3.0))
            tf_img2 = tx_img2.text_frame
            tf_img2.text = f"CÓDIGO INTERNO\n{codigo}"
                
            tf_img2.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in tf_img2.paragraphs:
                p.font.name = 'Calibri'
                p.font.size = Pt(14)
                p.font.bold = True
                p.alignment = PP_ALIGN.CENTER
                
            # Borda para o placeholder
            tx_img2.line.color.rgb = RGBColor(100, 100, 100)
            tx_img2.line.width = Pt(1)
            tx_img2.fill.solid()
            tx_img2.fill.fore_color.rgb = RGBColor(245, 245, 245)
            
            print(f"  [INFO] Adicionados placeholders - Ref: {ref if ref else 'N/A'} | Código Interno: {codigo}")
                        
    except Exception as e:
        print(f"  [ERRO] Erro ao adicionar imagens: {str(e)}")


def adicionar_especificacoes_sharepoint(slide, word_files, sp_client, produto=None):
    """Adiciona especificações técnicas do SharePoint ao slide"""
    try:
        if not word_files:
            # Se não encontrar arquivo Word, adicionar informações do produto
            if produto:
                codigo = produto.get('codigo', '')
                ref = produto.get('ref', '')
                descricao = produto.get('descricao', '')
                marca = produto.get('marca', '')
                
                # Priorizar referência sobre código interno
                if ref and ref != codigo:
                    texto_fallback = f"REFERÊNCIA: {ref}"
                    texto_fallback += f"\nCÓDIGO INTERNO: {codigo}"
                else:
                    texto_fallback = f"CÓDIGO ORÇAMENTO: {codigo}"
                    
                if marca:
                    texto_fallback += f"\nMARCA: {marca}"
                if descricao:
                    texto_fallback += f"\nDESCRIÇÃO: {descricao}"
                    
                texto_fallback += "\n\n[Ficha técnica não encontrada no SharePoint]"
                
                # Adicionar texto de fallback ao slide
                tx_txt = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(7.27), Inches(5.0))
                tf_txt = tx_txt.text_frame
                tf_txt.word_wrap = True
                tf_txt.text = texto_fallback
                
                for p in tf_txt.paragraphs:
                    p.font.name = 'Calibri'
                    p.font.size = Pt(14)
                    p.font.bold = True if "CÓDIGO ORÇAMENTO" in p.text else False
                    
                print(f"  [INFO] Adicionado código de orçamento no slide: {ref if ref and ref != codigo else codigo}")
            return
        
        # Extrair texto do primeiro arquivo Word
        download_url = word_files[0]['download_url']
        texto_specs = sp_client.get_word_text(download_url)
        
        # Adicionar texto ao slide
        tx_txt = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(7.27), Inches(5.0))
        tf_txt = tx_txt.text_frame
        tf_txt.word_wrap = True
        tf_txt.text = texto_specs
        
        for p in tf_txt.paragraphs:
            p.font.name = 'Calibri'
            p.font.size = Pt(14)
            
    except Exception as e:
        print(f"  [ERRO] Erro ao adicionar especificações: {str(e)}")


def sanitizar_texto_word(texto_word: str, componentes: List[Dict], mandante: Optional[Dict] = None) -> str:
    """
    Sanitiza texto extraído do Word removendo informações técnicas que conflitam com XML.
    
    Regras de precedência:
    - XML sempre tem prioridade sobre Word para especificações técnicas
    - Remove parágrafos do Word que mencionam componentes presentes no XML
    - Mantém informações gerais, dimensões, acabamentos, etc.
    
    Args:
        texto_word: Texto extraído do arquivo Word do SharePoint
        componentes: Lista de componentes do grupo (vindos do XML)
        mandante: Produto mandante (opcional, para contexto adicional)
    
    Returns:
        Texto sanitizado (sem conflitos técnicos)
    """
    try:
        if not componentes:
            # Sem componentes, não há conflito - retorna texto original
            return texto_word
        
        # Mapear categorias presentes nos componentes XML
        categorias_xml = set()
        for comp in componentes:
            categoria = comp.get('categoria', '').lower()
            if categoria:
                categorias_xml.add(categoria)
        
        # Palavras-chave mapeadas por categoria
        palavras_por_categoria = {
            'lâmpada': ['lâmpada', 'lampada', 'bulb', 'incandescente', 'halógena', 'halogena'],
            'led': ['led', 'l.e.d', 'chip led', 'smd'],
            'fita led': ['fita led', 'fita de led', 'strip led', 'led strip'],
            'driver': ['driver', 'fonte', 'transformador', 'conversor'],
            'reator': ['reator', 'ballast'],
            'acessório': ['acessório', 'acessorio', 'suporte', 'conector', 'fixação', 'fixacao'],
            'fonte': ['fonte', 'power supply', 'alimentação', 'alimentacao'],
        }
        
        # Palavras-chave técnicas gerais (sempre verificar)
        palavras_tecnicas_gerais = [
            'potência', 'potencia', 'watts', 'w/', 'watt',
            'lúmens', 'lumens', 'lm', 'fluxo luminoso',
            'temperatura de cor', 'kelvin', 'k', '2700k', '3000k', '4000k', '6500k',
            'voltagem', 'tensão', 'tensao', 'bivolt', '110v', '220v', '127v', '12v', '24v',
            'ip', 'grau de proteção', 'protecao', 'ip20', 'ip44', 'ip65', 'ip67',
            'irc', 'cri', 'índice de reprodução', 'indice de reproducao',
        ]
        
        # Identificar palavras-chave a filtrar baseado nas categorias XML
        palavras_filtrar = set()
        for categoria in categorias_xml:
            for key, palavras in palavras_por_categoria.items():
                if key in categoria:
                    palavras_filtrar.update(palavras)
        
        # Se há componentes técnicos (LED, Lâmpada, etc.), filtrar também especificações gerais
        if any(cat in categorias_xml for cat in ['led', 'fita led', 'lâmpada', 'lampada', 'driver', 'fonte']):
            palavras_filtrar.update(palavras_tecnicas_gerais)
        
        if not palavras_filtrar:
            # Sem palavras para filtrar, retorna texto original
            return texto_word
        
        # Processar texto linha por linha
        linhas_originais = texto_word.split('\n')
        linhas_sanitizadas = []
        linhas_removidas = 0
        
        for linha in linhas_originais:
            linha_lower = linha.lower()
            
            # Verificar se linha contém palavra-chave conflitante
            tem_conflito = any(palavra in linha_lower for palavra in palavras_filtrar)
            
            if tem_conflito:
                linhas_removidas += 1
                print(f"  [SANITIZACAO] Removida linha: {linha[:60]}{'...' if len(linha) > 60 else ''}")
            else:
                linhas_sanitizadas.append(linha)
        
        texto_sanitizado = '\n'.join(linhas_sanitizadas)
        
        if linhas_removidas > 0:
            print(f"  [SANITIZACAO] {linhas_removidas} linha(s) técnica(s) removida(s) para evitar conflito com XML")
            print(f"  [SANITIZACAO] Categorias XML detectadas: {', '.join(categorias_xml)}")
        
        return texto_sanitizado
        
    except Exception as e:
        print(f"  [ERRO] Erro ao sanitizar texto Word: {str(e)}")
        return texto_word  # Em caso de erro, retorna texto original


def adicionar_especificacoes_grupo(slide, word_files, sp_client, mandante, componentes):
    """
    Adiciona especificações técnicas do grupo ao slide
    Concatena informações do mandante + componentes
    
    Args:
        slide: Slide do PowerPoint
        word_files: Arquivos Word do mandante
        sp_client: Cliente SharePoint
        mandante: Produto mandante
        componentes: Lista de produtos componentes
    """
    try:
        # Texto base do mandante
        texto_specs = ""
        
        if word_files:
            # Extrair texto do primeiro arquivo Word (mandante)
            download_url = word_files[0]['download_url']
            
            # Usar filtro técnico se houver componentes no XML
            filtrar_tecnico = len(componentes) > 0
            texto_word_bruto = sp_client.get_word_text(download_url, filtrar_tecnico=filtrar_tecnico)
            
            # Aplicar sanitização adicional baseada nos componentes específicos
            if componentes:
                print(f"  [PRECEDENCIA] Aplicando sanitização - XML tem prioridade sobre Word")
                texto_specs = sanitizar_texto_word(texto_word_bruto, componentes, mandante)
            else:
                texto_specs = texto_word_bruto
        else:
            # Fallback: informações do mandante
            codigo = mandante.get('codigo', '')
            ref = mandante.get('ref', '')
            descricao = mandante.get('descricao', '')
            marca = mandante.get('marca', '')
            categoria = mandante.get('categoria', '')
            
            if ref and ref != codigo:
                texto_specs = f"REFERÊNCIA: {ref}\n"
                texto_specs += f"CÓDIGO INTERNO: {codigo}\n"
            else:
                texto_specs = f"CÓDIGO ORÇAMENTO: {codigo}\n"
            
            if marca:
                texto_specs += f"MARCA: {marca}\n"
            if categoria:
                texto_specs += f"CATEGORIA: {categoria}\n"
            if descricao:
                texto_specs += f"DESCRIÇÃO: {descricao}\n"
        
        # Adicionar informações dos componentes (PRECEDÊNCIA: XML > Word)
        if componentes:
            print(f"  [COMPONENTES] Adicionando {len(componentes)} componente(s) ao slide (dados do XML)")
            
            # Separador visual
            texto_specs += "\n" + "="*60 + "\n"
            texto_specs += "ESPECIFICAÇÕES TÉCNICAS (do Projeto):\n"
            texto_specs += "="*60 + "\n\n"
            
            for idx, comp in enumerate(componentes, 1):
                codigo_comp = comp.get('codigo', '')
                ref_comp = comp.get('ref', '')
                desc_comp = comp.get('descricao', '')
                categoria_comp = comp.get('categoria', '')
                qtd_comp = comp.get('quantidade', '1')
                unidade_comp = comp.get('unidade', 'UN')
                
                # Montar texto do componente com mais detalhes
                texto_comp = f"[{idx}] "
                
                if categoria_comp:
                    texto_comp += f"{categoria_comp.upper()} - "
                
                if ref_comp and ref_comp != codigo_comp:
                    texto_comp += f"Ref: {ref_comp}"
                    if codigo_comp:
                        texto_comp += f" (Cód. Interno: {codigo_comp})"
                else:
                    texto_comp += f"Código: {codigo_comp}"
                
                if qtd_comp != '1':
                    texto_comp += f" | Qtd: {qtd_comp} {unidade_comp}"
                
                if desc_comp:
                    # Limitar descrição a 100 caracteres por componente
                    desc_curta = desc_comp[:100] + "..." if len(desc_comp) > 100 else desc_comp
                    texto_comp += f"\n    → {desc_curta}"
                
                texto_specs += texto_comp + "\n\n"
                
                print(f"    [{idx}] {ref_comp or codigo_comp}: {categoria_comp or 'SEM CATEGORIA'} (Qtd: {qtd_comp})")
            
            # Nota sobre precedência
            texto_specs += "\n" + "-"*60 + "\n"
            texto_specs += "* Dados extraídos do orçamento XML (versão mais recente)\n"
        
        # Limitar tamanho do texto (máximo ~1000 caracteres para caber no slide)
        if len(texto_specs) > 1000:
            texto_specs = texto_specs[:1000] + "...\n\n[Texto truncado]"
        
        # Adicionar texto ao slide
        tx_txt = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(7.27), Inches(5.0))
        tf_txt = tx_txt.text_frame
        tf_txt.word_wrap = True
        tf_txt.text = texto_specs
        
        # Formatação
        for p in tf_txt.paragraphs:
            p.font.name = 'Calibri'
            p.font.size = Pt(11)  # Fonte menor para caber mais texto
            
            # Destacar títulos
            if any(keyword in p.text for keyword in ['REFERÊNCIA:', 'CÓDIGO', 'COMPONENTES ADICIONAIS:', '====']):
                p.font.bold = True
                
    except Exception as e:
        print(f"  [ERRO] Erro ao adicionar especificações do grupo: {str(e)}")


def criar_slide_bula(prs, bula_file, lnum, sp_client):
    """Cria slide dedicado para bula de instalação"""
    try:
        slide_bula = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Número LXX no canto superior direito em preto
        tx_lnum_bula = slide_bula.shapes.add_textbox(Inches(6.8), Inches(0.3), Inches(1), Inches(0.5))
        tf_lnum_bula = tx_lnum_bula.text_frame
        tf_lnum_bula.text = f"L{lnum}"
        tf_lnum_bula.paragraphs[0].font.name = 'Calibri'
        tf_lnum_bula.paragraphs[0].font.size = Pt(36)
        tf_lnum_bula.paragraphs[0].font.bold = True
        tf_lnum_bula.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        tf_lnum_bula.paragraphs[0].alignment = PP_ALIGN.RIGHT
        
        # Detectar extensão do arquivo de bula
        nome_bula = bula_file.get('name', '')
        extensao = os.path.splitext(nome_bula)[1].lower()
        if extensao not in ['.jpg', '.jpeg', '.png', '.pdf']:
            extensao = '.png'  # Fallback
        
        # Baixar e adicionar imagem da bula
        img_content = sp_client.download_file_content(bula_file['download_url'])
        
        with tempfile.NamedTemporaryFile(suffix=extensao, delete=False) as temp_img:
            temp_img.write(img_content)
            temp_path = temp_img.name
        
        try:
            slide_bula.shapes.add_picture(temp_path, Inches(0.5), Inches(1.0), width=Inches(7.2))
        finally:
            os.unlink(temp_path)
            
    except Exception as e:
        print(f"  [ERRO] Erro ao criar slide bula: {str(e)}")


# Função main para compatibilidade com sistema existente
def main():
    """Função principal - chama geração com SharePoint"""
    try:
        # TODO: Estes caminhos serão passados dinamicamente pelo app.py
        xml_path = xml_orcamento  
        excel_path = excel_master  
        
        return gerar_powerpoint_sharepoint(xml_path, excel_path)
        
    except Exception as e:
        print(f"[ERRO] Erro na execução principal: {str(e)}")
        raise


# Extrai produtos (substituindo código antigo do PDF)
produtos = []  # Será preenchido pela função gerar_powerpoint_sharepoint()

prs = Presentation()
prs.slide_width = a4_width
prs.slide_height = a4_height

# Header barra preta pequena (apenas para o título, mais abaixo)
def add_header(slide):
    # Barra preta menor - largura reduzida e altura apenas para o título
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(0.5), Inches(3.5), Inches(0.5))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0, 0, 0)
    txbox = slide.shapes.add_textbox(Inches(0.4), Inches(0.6), Inches(3.3), Inches(0.3))
    tf = txbox.text_frame
    tf.text = "LUCENERA - ATELIER DA LUZ"
    tf.paragraphs[0].font.name = 'Calibri'
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    tf.paragraphs[0].font.bold = True


# Executar quando chamado diretamente
if __name__ == "__main__":
    try:
        resultado = main()
        print(f"🎉 Processamento concluído: {resultado}")
    except Exception as e:
        print(f"💥 Falha na execução: {str(e)}")