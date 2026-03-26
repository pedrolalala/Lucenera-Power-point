import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
import pdfplumber
from docx import Document
from PIL import Image
import io

pasta_separada = r"C:\Users\pedro\OneDrive\Desktop\lucenera\INTERLIGHT_SEPARADOS"
pdf_orcamento = r"C:\Users\pedro\OneDrive\Desktop\lucenera\OR_0006316 FELIPE E DENISE PAISAGISMO b.pdf"
logo_path = r"C:\script python\script python power point\logo.png"
ppt_saida = r"C:\Users\pedro\OneDrive\Desktop\lucenera\Fichas_técnicas_v10.pptx"

img_exts = {'.jpg', '.jpeg', '.png', '.gif', '.bmp'}

a4_width = Inches(8.27)
a4_height = Inches(11.69)

# Extrai produtos
produtos = []  # Mudado de dict para lista para suportar códigos duplicados
with pdfplumber.open(pdf_orcamento) as pdf:
    text = pdf.pages[0].extract_text()
    matches = re.finditer(r'L0(\d)\s+(\d+)\s+([A-Z0-9\-]+).*?R\$\s*([\d.,]+)', text, re.DOTALL)
    for match in matches:
        lnum = match.group(1)
        codigo = match.group(2)
        ref = match.group(3)
        preco = match.group(4)
        produtos.append({"lnum": lnum, "codigo": codigo, "ref": ref, "preco": preco})

prs = Presentation()
prs.slide_width = a4_width
prs.slide_height = a4_height

# Header barra preta pequena (apenas para o título, mais abaixo)
def add_header(slide):
    # Barra preta menor - largura reduzida e altura apenas para o título
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(0.5), Inches(3.5), Inches(0.5))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0, 0, 0)
    txbox = slide.shapes.add_textbox(Inches(0.4), Inches(0.6), Inches(3.3), Inches(0.3))
    tf = txbox.text_frame
    tf.text = "LUCENERA - ATELIER DA LUZ"
    tf.paragraphs[0].font.name = 'Calibri'
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    tf.paragraphs[0].font.bold = True



# Capa (logo central)
slide_capa = prs.slides.add_slide(prs.slide_layouts[6])
if os.path.exists(logo_path):
    # Logo centralizada no meio da página
    logo_width = 3
    logo_height = 3
    logo_left = (8.27 - logo_width) / 2
    logo_top = (11.69 - logo_height) / 2
    slide_capa.shapes.add_picture(logo_path, Inches(logo_left), Inches(logo_top), width=Inches(logo_width))

# Por produto
for info in produtos:  # Mudado para iterar sobre a lista
    codigo = info['codigo']  # Extrair código do dict
    # Ficha (header + título pequeno esq + L0X grande dir + 2 fotos centro + specs baixo)
    slide_ficha = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide_ficha)
    
    # Título pequeno em preto abaixo do cabeçalho
    tx_title = slide_ficha.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(3), Inches(0.4))
    tf_title = tx_title.text_frame
    tf_title.text = "LUMINÁRIA DE PROJETO"
    tf_title.paragraphs[0].font.name = 'Calibri'
    tf_title.paragraphs[0].font.size = Pt(12)
    tf_title.paragraphs[0].font.bold = False
    tf_title.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Preto
    
    # Número L0X grande em preto no canto superior direito
    tx_lnum = slide_ficha.shapes.add_textbox(Inches(6.8), Inches(1.2), Inches(1), Inches(0.5))
    tf_lnum = tx_lnum.text_frame
    tf_lnum.text = f"L0{info['lnum']}"
    tf_lnum.paragraphs[0].font.name = 'Calibri'
    tf_lnum.paragraphs[0].font.size = Pt(36)
    tf_lnum.paragraphs[0].font.bold = True
    tf_lnum.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Preto
    tf_lnum.paragraphs[0].alignment = PP_ALIGN.RIGHT
    
    # Linha contínua ANTES das fotos
    line1 = slide_ficha.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.5), Inches(2.0), Inches(7.77), Inches(2.0))
    line1.line.color.rgb = RGBColor(0, 0, 0)
    
    # Duas fotos lado a lado no centro (luminária + ambiente)
    foto_produto = None
    foto_ambiente = None
    docs = [f for f in os.listdir(pasta_separada) if codigo in f.upper() and f.endswith('.docx')]
    imgs = [f for f in os.listdir(pasta_separada) if codigo in f.upper() and os.path.splitext(f)[1].lower() in img_exts]
    
    # Extrair foto do produto do .docx
    if docs:
        doc_path = os.path.join(pasta_separada, docs[0])
        doc = Document(doc_path)
        img_count = 0
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                img_data = rel.target_part.blob
                img_stream = io.BytesIO(img_data)
                img = Image.open(img_stream)
                temp_path = f"temp_{codigo}_{img_count}.png"
                img.save(temp_path)
                if img_count == 0:
                    foto_produto = temp_path
                elif img_count == 1:
                    foto_ambiente = temp_path
                img_count += 1
                if img_count >= 2:
                    break
    
    # Se não tiver no .docx, tentar imagens separadas
    if not foto_produto and imgs:
        foto_produto = os.path.join(pasta_separada, imgs[0])
    
    # Adicionar fotos lado a lado centralizadas (técnica menor esq, produto maior dir)
    if foto_produto:
        slide_ficha.shapes.add_picture(foto_produto, Inches(1.0), Inches(2.3), height=Inches(2.5))
        if foto_produto.startswith("temp_"):
            os.remove(foto_produto)
    
    if foto_ambiente:
        slide_ficha.shapes.add_picture(foto_ambiente, Inches(4.2), Inches(2.3), height=Inches(3.0))
        if foto_ambiente.startswith("temp_"):
            os.remove(foto_ambiente)
    
    # Linha contínua DEPOIS das fotos
    line2 = slide_ficha.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.5), Inches(5.5), Inches(7.77), Inches(5.5))
    line2.line.color.rgb = RGBColor(0, 0, 0)
    
    # Specs embaixo da segunda linha
    tx_txt = slide_ficha.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(7.27), Inches(5.0))
    tf_txt = tx_txt.text_frame
    tf_txt.word_wrap = True
    if docs:
        doc = Document(os.path.join(pasta_separada, docs[0]))
        # Filtrar linhas vazias, linhas com apenas underscores/traços, e títulos duplicados
        texto_filtrado = []
        for p in doc.paragraphs:
            texto_linha = p.text.strip()
            # Ignorar linhas vazias, separadores e títulos que já estão no cabeçalho
            if (texto_linha and 
                not all(c in '_-—–=' for c in texto_linha) and
                'LUCENERA' not in texto_linha.upper() and
                'ATELIER DA LUZ' not in texto_linha.upper() and
                'LUMINÁRIA DE PROJETO' not in texto_linha.upper()):
                texto_filtrado.append(texto_linha)
        texto = '\n'.join(texto_filtrado)
        tf_txt.text = texto[:700] + "..."
    for p in tf_txt.paragraphs:
        p.font.name = 'Calibri'
        p.font.size = Pt(14)
    

    # Bula (sem cabeçalho, apenas imagem full) - só cria se encontrar a bula
    bulas = [f for f in os.listdir(pasta_separada) if ('BULA' in f.upper() or 'INSTALACAO' in f.upper()) and codigo in f.upper() and os.path.splitext(f)[1].lower() in img_exts]
    if bulas:
        slide_bula = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Número L0X no canto superior direito em preto
        tx_lnum_bula = slide_bula.shapes.add_textbox(Inches(6.8), Inches(0.3), Inches(1), Inches(0.5))
        tf_lnum_bula = tx_lnum_bula.text_frame
        tf_lnum_bula.text = f"L0{info['lnum']}"
        tf_lnum_bula.paragraphs[0].font.name = 'Calibri'
        tf_lnum_bula.paragraphs[0].font.size = Pt(36)
        tf_lnum_bula.paragraphs[0].font.bold = True
        tf_lnum_bula.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        tf_lnum_bula.paragraphs[0].alignment = PP_ALIGN.RIGHT
        
        # Imagem da bula ocupando toda a página
        img_path = os.path.join(pasta_separada, bulas[0])
        slide_bula.shapes.add_picture(img_path, Inches(0.5), Inches(1.0), width=Inches(7.2), height=Inches(9.5))
       

prs.save(ppt_saida)
print(f"✅ PPT v10 fixado: {ppt_saida} (import MSO_SHAPE + layout exato)!")